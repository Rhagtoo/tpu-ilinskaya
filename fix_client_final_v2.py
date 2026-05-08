#!/usr/bin/env python3
"""
Финальный патч презентации ТПУ Ильинская.

Назначение:
    Применяет последние клиентские правки к файлу TPU_Ilinskaya_client_final.pptx
    и сохраняет финальную версию TPU_Ilinskaya_client_final_v2.pptx.

Что делает:
    - убирает крупную надпись «ФИНАЛЬНАЯ ВЕРСИЯ» со слайда структуры;
    - пересобирает слайды-разделители в корпоративной стилистике bnMAP;
    - удаляет текущие слайды 12 и 14;
    - исправляет строку площади на слайде 4: «206 700 м²» в одну строку;
    - убирает ошибочный маркер «МКАД» с транспортной карты на слайде 5;
    - пересобирает слайд позиционирования конкурентов;
    - оформляет последний слайд в стиле bnMAP.

Запуск из корня проекта:
    python scripts/fix_client_final_v2.py \
        --input TPU_Ilinskaya_client_final.pptx \
        --output TPU_Ilinskaya_client_final_v2.pptx

Зависимости:
    pip install python-pptx pillow opencv-python numpy
"""

from __future__ import annotations

import argparse
import os
import zipfile
from pathlib import Path

import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# =============================================================================
# DESIGN
# =============================================================================

W, H = 1600, 900
BLUE = (43, 82, 197)
DARK = (13, 33, 55)
ORANGE = (245, 156, 38)
WHITE = (255, 255, 255)
TEXT = (25, 25, 25)
MUTED = (90, 100, 115)


# =============================================================================
# HELPERS: raster assets
# =============================================================================

def font(size: int = 40, bold: bool = False):
    candidates = []
    if bold:
        candidates += [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/Library/Fonts/Arial Bold.ttf",
            "C:/Windows/Fonts/arialbd.ttf",
        ]
    candidates += [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/Library/Fonts/Arial.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for p in candidates:
        if os.path.exists(p):
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()


def draw_logo(d: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0, color=WHITE) -> None:
    # Текстовая имитация логотипа без вшивания шрифтов/брендовых файлов.
    d.text((x, y), "bnMAP.pro", font=font(int(42 * scale), True), fill=color)


def make_cover(path: Path) -> None:
    img = Image.new("RGB", (W, H), BLUE)
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)

    for r, alpha in [(680, 60), (520, 45), (360, 35)]:
        od.ellipse((W - r // 2, -r // 2, W + r, r), fill=(0, 190, 205, alpha))
    od.polygon([(0, H), (0, 620), (510, H)], fill=(10, 34, 83, 80))
    od.polygon([(900, H), (W, H), (W, 540)], fill=(18, 51, 114, 75))
    for i in range(14):
        x0 = -120 + i * 150
        od.line([(x0, 760), (x0 + 520, 440), (x0 + 900, 520)], fill=(255, 255, 255, 20), width=2)

    img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    d = ImageDraw.Draw(img)

    draw_logo(d, W - 330, 54, 1.0)
    d.text((92, 92), "Цифровой сервис мониторинга\nи анализа рынка новостроек", font=font(24), fill=WHITE)
    d.line((92, 206, 560, 206), fill=ORANGE, width=4)
    d.text((92, 300), "АНАЛИТИЧЕСКИЙ ОТЧЁТ", font=font(32), fill=(220, 232, 255))
    d.text((92, 365), "ТПУ ИЛЬИНСКАЯ", font=font(72, True), fill=WHITE)
    d.text((92, 470), "г. Москва, Новорижское шоссе, 22 км", font=font(30), fill=(235, 245, 255))

    d.rounded_rectangle((980, 315, 1460, 610), radius=34, fill=(255, 255, 255), outline=(255, 255, 255), width=2)
    # Полупрозрачность rectangle через overlay не нужна: оставляем чистый корпоративный блок.
    d.rounded_rectangle((980, 315, 1460, 610), radius=34, fill=(75, 116, 218), outline=(135, 169, 235), width=2)
    d.text((1028, 365), "ЗАПАД", font=font(46, True), fill=WHITE)
    d.text((1028, 438), "Новорижское шоссе\nст. Ильинская", font=font(28), fill=(232, 244, 255), spacing=8)
    d.text((92, 805), "Май 2026 г.  ·  Конфиденциально", font=font(24), fill=(220, 235, 255))

    img.save(path, quality=95)


def make_section(path: Path, num: int, subtitle: str) -> None:
    img = Image.new("RGB", (W, H), WHITE)
    d = ImageDraw.Draw(img)

    d.rectangle((0, 0, W, 112), fill=BLUE)
    d.text((84, 32), f"РАЗДЕЛ {num}", font=font(46), fill=WHITE)
    draw_logo(d, W - 270, 32, 0.82)

    grad = Image.new("RGBA", (W, H), (255, 255, 255, 0))
    gd = ImageDraw.Draw(grad)
    for x in range(W):
        t = max(0, min(1, (x - 650) / (W - 650)))
        if t > 0:
            gd.line((x, 112, x, H), fill=(145, 185, 245, int(210 * t)), width=1)
    gd.ellipse((900, 120, 1900, 940), fill=(70, 116, 217, 115))
    gd.ellipse((700, 250, 1550, 1100), fill=(115, 164, 235, 70))

    img = Image.alpha_composite(img.convert("RGBA"), grad).convert("RGB")
    d = ImageDraw.Draw(img)

    d.text((122, 285), subtitle, font=font(31, True), fill=(0, 0, 0))
    d.line((122, 365, 720, 365), fill=ORANGE, width=3)

    # Простая корпоративная изометрическая иллюстрация.
    cx, cy = 1220, 560
    d.polygon([(cx - 210, cy + 60), (cx, cy - 50), (cx + 235, cy + 65), (cx + 20, cy + 190)], fill=(246, 249, 255), outline=(222, 232, 248))
    d.rounded_rectangle((cx - 10, cy - 190, cx + 130, cy + 50), radius=18, fill=BLUE, outline=(88, 125, 220), width=2)
    d.rounded_rectangle((cx + 20, cy - 225, cx + 155, cy - 12), radius=8, fill=(236, 243, 255), outline=(210, 223, 242), width=2)
    d.rectangle((cx + 55, cy - 155, cx + 90, cy - 95), fill=(88, 132, 229))
    d.line((cx + 40, cy - 68, cx + 130, cy - 68), fill=(203, 215, 235), width=4)
    d.line((cx + 40, cy - 43, cx + 115, cy - 43), fill=(203, 215, 235), width=4)

    d.polygon([(900, 310), (990, 265), (1080, 310), (990, 360)], fill=(245, 249, 255), outline=(222, 232, 248))
    d.line((952, 270, 952, 335), fill=BLUE, width=8)
    d.line((1040, 290, 1040, 355), fill=BLUE, width=8)
    d.line((955, 275, 1040, 295), fill=BLUE, width=7)
    d.line([(968, 320), (990, 306), (1012, 313), (1036, 288)], fill=(255, 90, 65), width=4)

    img.save(path, quality=95)


def make_contact(path: Path) -> None:
    img = Image.new("RGB", (W, H), BLUE)
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.ellipse((760, -140, 1920, 980), fill=(0, 180, 190, 105))
    od.polygon([(0, H), (0, 630), (420, H)], fill=(13, 33, 55, 100))
    img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    d = ImageDraw.Draw(img)

    draw_logo(d, 94, 66, 1.0)
    d.text((94, 128), "Цифровой сервис мониторинга и анализа рынка новостроек", font=font(26), fill=(220, 235, 255))
    d.line((94, 185, 590, 185), fill=ORANGE, width=4)
    d.text((94, 300), "КОНТАКТЫ", font=font(62, True), fill=WHITE)

    d.rounded_rectangle((850, 250, 1460, 640), radius=28, fill=(255, 255, 255), outline=(255, 255, 255), width=2)
    x, y = 910, 315
    d.text((x, y), "Телефон", font=font(26, True), fill=DARK)
    d.text((x, y + 44), "+7 903 149 66 50", font=font(30), fill=TEXT)
    d.text((x, y + 118), "E-mail", font=font(26, True), fill=DARK)
    d.text((x, y + 162), "smart@itlabrealty.ru", font=font(30), fill=TEXT)
    d.text((x, y + 236), "Сайт", font=font(26, True), fill=DARK)
    d.text((x, y + 280), "bnmap.pro", font=font(30), fill=TEXT)

    d.text((94, 740), "Аналитический отчёт по проекту ТПУ Ильинская\nМай 2026 г.  ·  Конфиденциально", font=font(26), fill=(225, 238, 255), spacing=8)
    img.save(path, quality=95)


def make_assets(asset_dir: Path) -> tuple[Path, Path, dict[int, Path]]:
    asset_dir.mkdir(parents=True, exist_ok=True)
    cover = asset_dir / "cover.png"
    contact = asset_dir / "contact.png"
    sections = {
        1: asset_dir / "section1.png",
        2: asset_dir / "section2.png",
        3: asset_dir / "section3.png",
    }
    make_cover(cover)
    make_contact(contact)
    make_section(sections[1], 1, "Анализ местоположения Проекта")
    make_section(sections[2], 2, "Обзор проектов конкурентного окружения в локации")
    make_section(sections[3], 3, "Ценообразование и стратегия реализации")
    return cover, contact, sections


# =============================================================================
# HELPERS: pptx
# =============================================================================

def remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def clear_slide(slide) -> None:
    for sh in list(slide.shapes):
        remove_shape(sh)


def delete_slide(prs: Presentation, index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    r_id = xml_slides[index].get(qn("r:id"))
    prs.part.drop_rel(r_id)
    xml_slides.remove(xml_slides[index])


def set_text(shape, text: str, size=None, bold=None, color=None) -> None:
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = False if "м²" in text else tf.word_wrap
    for p in tf.paragraphs:
        for r in p.runs:
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color
            r.font.name = "Calibri"


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=RGBColor(25, 25, 25), align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Calibri"
    r.font.color.rgb = color
    return tb


def rebuild_positioning(slide, prs) -> None:
    clear_slide(slide)
    sw = prs.slide_width
    navy = RGBColor(43, 82, 197)
    dark = RGBColor(25, 25, 25)
    white_c = RGBColor(255, 255, 255)
    light_bg = RGBColor(245, 248, 253)
    line = RGBColor(220, 226, 235)

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.58))
    rect.fill.solid()
    rect.fill.fore_color.rgb = navy
    rect.line.fill.background()

    add_textbox(slide, 8.05, 0.13, 1.55, 0.28, "bnMAP.pro", size=17, bold=True, color=white_c, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.55, 0.92, 8.9, 0.38, "КОНКУРЕНТНАЯ ВЫБОРКА: ПОЗИЦИОНИРОВАНИЕ, ПРЕИМУЩЕСТВА И НЕДОСТАТКИ", size=17, bold=True, color=dark)

    x0, y0 = 0.55, 1.47
    col1, col2, col3 = 3.05, 2.8, 3.05
    head_h, row_h = 0.42, 0.78

    x = x0
    for header, width in [("Проект / позиционирование", col1), ("Преимущества", col2), ("Недостатки", col3)]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0), Inches(width), Inches(head_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = navy
        shape.line.color.rgb = navy
        add_textbox(slide, x + 0.08, y0 + 0.09, width - 0.16, 0.22, header, size=10.5, bold=True, color=white_c)
        x += width

    rows = [
        (
            "СберСити",
            "Первый инновационный смарт-район России на берегу Москвы-реки.",
            "Уникальный масштаб и инфраструктура; бренд Сбербанка; ИИ в ЖКХ; 7 школ, 13 детсадов.",
            "Высокие цены; значительный объём предложения; конкуренция с проектом по статусу и масштабу.",
        ),
        (
            "Станиславский",
            "Малоэтажный жилой квартал в датско-голландском стиле в 10 мин от Москвы.",
            "Камерность, малоэтажность; авторская архитектура; панорамные террасы; варианты отделки.",
            "Новый девелопер; близость к Новорижскому шоссе; нет метро в пешей доступности.",
        ),
        (
            "Мыс",
            "Загородный премиум-проект 200 га в Ликино: урбан-блоки, клубные дома, таунхаусы, коттеджи.",
            "Природное расположение; разнообразие форматов; инфраструктура курортного уровня.",
            "Удалённость от метро; более высокий ценовой уровень; длительный срок реализации.",
        ),
        (
            "Строгино 360",
            "Флагман ПИК в экологичном районе Строгино с пешеходным бульваром.",
            "Метро в пешей доступности; бульвар; природный парк рядом; сильная транспортная доступность.",
            "Высокая плотность застройки; самый высокий бюджет покупки среди комфорт-класса.",
        ),
    ]

    y = y0 + head_h
    for ri, row in enumerate(rows):
        bg = RGBColor(255, 255, 255) if ri % 2 == 0 else light_bg
        x = x0
        for width in [col1, col2, col3]:
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(row_h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = bg
            shp.line.color.rgb = line
            shp.line.width = Pt(0.6)
            x += width

        add_textbox(slide, x0 + 0.12, y + 0.10, col1 - 0.24, 0.18, row[0], size=11.5, bold=True, color=navy)
        add_textbox(slide, x0 + 0.12, y + 0.32, col1 - 0.24, 0.34, row[1], size=8.5, color=dark)
        add_textbox(slide, x0 + col1 + 0.12, y + 0.12, col2 - 0.24, 0.56, row[2], size=8.7, color=dark)
        add_textbox(slide, x0 + col1 + col2 + 0.12, y + 0.12, col3 - 0.24, 0.56, row[3], size=8.7, color=dark)
        y += row_h


def extract_media_image(input_pptx: Path, filename: str, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    out = output_dir / filename
    with zipfile.ZipFile(input_pptx, "r") as zf:
        candidates = [name for name in zf.namelist() if name.endswith(f"ppt/media/{filename}")]
        if not candidates:
            raise FileNotFoundError(f"Не найдено ppt/media/{filename} внутри {input_pptx}")
        out.write_bytes(zf.read(candidates[0]))
    return out


def patch_transport_image(src: Path, dst: Path) -> None:
    im = Image.open(src).convert("RGB")
    arr = np.array(im)
    mask = np.zeros(arr.shape[:2], dtype=np.uint8)
    # Маска закрывает только ошибочный кастомный маркер/лейбл МКАД, не саму карту.
    mask[285:355, 650:780] = 255
    out = cv2.inpaint(arr, mask, 7, cv2.INPAINT_TELEA)
    Image.fromarray(out).save(dst, quality=95)


def is_picture_shape(shape) -> bool:
    # MSO_SHAPE_TYPE.PICTURE == 13; не импортируем enum ради совместимости.
    return getattr(shape, "shape_type", None) == 13


# =============================================================================
# MAIN
# =============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(description="Финальный патч презентации ТПУ Ильинская")
    parser.add_argument("--input", default="TPU_Ilinskaya_client_final.pptx", help="Входная презентация")
    parser.add_argument("--output", default="TPU_Ilinskaya_client_final_v2.pptx", help="Выходная презентация")
    parser.add_argument("--assets", default="final_assets", help="Папка для временных PNG-ассетов")
    args = parser.parse_args()

    input_pptx = Path(args.input).resolve()
    output_pptx = Path(args.output).resolve()
    asset_dir = Path(args.assets).resolve()

    if not input_pptx.exists():
        raise FileNotFoundError(f"Не найден входной файл: {input_pptx}")

    cover, contact, sections = make_assets(asset_dir)

    # Извлекаем исходную транспортную карту из PPTX, чтобы не требовать отдельную папку client_media.
    image7 = extract_media_image(input_pptx, "image7.png", asset_dir)
    transport_fixed = asset_dir / "transport_nomkad.png"
    patch_transport_image(image7, transport_fixed)

    prs = Presentation(str(input_pptx))

    # Удаляем текущие слайды 14 и 12, с конца, чтобы индексы не съехали.
    for idx in [13, 11]:
        if idx < len(prs.slides):
            delete_slide(prs, idx)

    # 1. Cover.
    clear_slide(prs.slides[0])
    prs.slides[0].shapes.add_picture(str(cover), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # 2. Structure: убрать правую финальную плашку и обновить диапазоны.
    s = prs.slides[1]
    for sh in list(s.shapes):
        txt = getattr(sh, "text", "") if hasattr(sh, "text") else ""
        if ("ФИНАЛЬНАЯ" in txt) or ("инвесторский" in txt) or (sh.left > Inches(6.7) and sh.top > Inches(1.2)):
            remove_shape(sh)

    for sh in s.shapes:
        if hasattr(sh, "text"):
            if sh.text.strip() == "Слайды 3–12":
                set_text(sh, "Слайды 3–11", size=12, color=RGBColor(80, 90, 110))
            elif sh.text.strip() == "Слайды 13–15":
                set_text(sh, "Слайды 12–13", size=12, color=RGBColor(80, 90, 110))
            elif sh.text.strip() == "Слайды 16–18":
                set_text(sh, "Слайды 14–16", size=12, color=RGBColor(80, 90, 110))

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.65), Inches(2.05), Inches(2.25))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 248, 253)
    box.line.color.rgb = RGBColor(220, 226, 235)
    add_textbox(s, 7.45, 2.12, 1.55, 0.3, "ТПУ Ильинская", size=15, bold=True, color=RGBColor(43, 82, 197), align=PP_ALIGN.CENTER)
    add_textbox(s, 7.32, 2.52, 1.82, 0.7, "аналитический отчёт\nмай 2026", size=11, color=RGBColor(65, 75, 95), align=PP_ALIGN.CENTER)

    # Section slides.
    for slide in prs.slides:
        text = "\n".join([getattr(sh, "text", "") for sh in slide.shapes if hasattr(sh, "text")])
        if "РАЗДЕЛ 1" in text:
            clear_slide(slide)
            slide.shapes.add_picture(str(sections[1]), 0, 0, width=prs.slide_width, height=prs.slide_height)
        elif "РАЗДЕЛ 2" in text:
            clear_slide(slide)
            slide.shapes.add_picture(str(sections[2]), 0, 0, width=prs.slide_width, height=prs.slide_height)
        elif "РАЗДЕЛ 3" in text:
            clear_slide(slide)
            slide.shapes.add_picture(str(sections[3]), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Slide 4: площадь в одну строку.
    slide4 = prs.slides[3]
    for sh in slide4.shapes:
        if hasattr(sh, "text") and sh.text.strip() == "206 700 м²":
            sh.left = Inches(8.05)
            sh.width = Inches(1.15)
            sh.height = Inches(0.28)
            sh.text_frame.word_wrap = False
            set_text(sh, "206 700 м²", size=13.5, bold=True, color=RGBColor(13, 33, 55))

    # Slide 5: заменить карту с ошибочным маркером МКАД на inpaint-версию.
    slide5 = prs.slides[4]
    for sh in list(slide5.shapes):
        if is_picture_shape(sh) and abs(sh.left - Inches(0.55)) < 10000 and abs(sh.top - Inches(0.98)) < 10000:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            remove_shape(sh)
            slide5.shapes.add_picture(str(transport_fixed), l, t, width=w, height=h)
            break

    # Rebuild positioning slide.
    for slide in prs.slides:
        text = "\n".join([getattr(sh, "text", "") for sh in slide.shapes if hasattr(sh, "text")])
        if "КОНКУРЕНТНАЯ ВЫБОРКА: ПОЗИЦИОНИРОВАНИЕ" in text:
            rebuild_positioning(slide, prs)
            break

    # Last slide: contact.
    last = prs.slides[-1]
    clear_slide(last)
    last.shapes.add_picture(str(contact), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(str(output_pptx))
    print(f"Saved: {output_pptx}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
