#!/usr/bin/env python3
"""
Безопасная сборка ТПУ Ильинская — final.

Архитектура:
- НЕ удаляет существующие слайды.
- НЕ пересобирает разделы 2/3.
- НЕ добавляет карту-подложку на слайд конкурентного окружения.
- Перерисовывает существующие слайды 4–8:
    4 — Об участке и местоположении
    5 — Транспортная доступность
    6 — Образование и ДОУ
    7 — Спорт и рекреация
    8 — Торговля
- Добавляет новый слайд 9 — Медицина.

Важно:
Для появления отдельного слайда 9 PowerPoint-файл всё равно должен получить новый slide_id.
Скрипт создаёт новый слайд штатным prs.slides.add_slide(), затем перемещает ТОЛЬКО этот новый slide_id
на позицию 9. Существующие slide parts/relationships не удаляются и не клонируются.

Перед запуском:
    python generate_maps_final.py
    python build_safe_final.py

Зависимости:
    pip install python-pptx pillow lxml
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# =============================================================================
# CONFIG
# =============================================================================

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
SOURCE_PPTX = os.path.join(PROJECT_DIR, "TPU_Ilinskaya_reworked_final.pptx")
OUTPUT_PPTX = os.path.join(PROJECT_DIR, "TPU_Ilinskaya_final_safe.pptx")
MAPS_DIR = os.path.join(PROJECT_DIR, "maps")

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.63)

# =============================================================================
# DESIGN SYSTEM
# =============================================================================

NAVY = RGBColor(0x0D, 0x21, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x19, 0x19, 0x19)
MUTED_TEXT = RGBColor(0x66, 0x77, 0x88)
PANEL_BG = RGBColor(0xF8, 0xF8, 0xF8)
PANEL_BG_BLUE = RGBColor(0xED, 0xF4, 0xFA)
ACCENT = RGBColor(0x0B, 0x5D, 0x8C)
FONT_FAMILY = "Calibri"

# =============================================================================
# DATA
# =============================================================================

@dataclass
class Scene:
    name: str
    title: str
    bullets: List[str]
    conclusion: str
    panel_blue: bool = False


SCENES: List[Scene] = [
    Scene(
        name="location",
        title="ОБ УЧАСТКЕ И МЕСТОПОЛОЖЕНИИ",
        bullets=[
            "Кадастровый номер: 50:11:0050603:423",
            "Площадь участка: ~20,6 га",
            "Форма участка: вытянутый линейный тип",
            "Север: Новорижское шоссе",
            "Юг: Москва-река",
            "Запад: Воронки / Архангельское",
            "Восток: Рублёво / МКАД направление",
        ],
        conclusion="Участок расположен в зоне активного развития западного направления Москвы.",
    ),
    Scene(
        name="transport",
        title="ТРАНСПОРТНАЯ ДОСТУПНОСТЬ",
        bullets=[
            "Расстояние до МКАД: ~5 км",
            "Москва-Сити: ~40–55 мин",
            "м. Строгино: ~20–30 мин",
            "ст. Ильинская: шаговая доступность",
            "Новорижское шоссе — основной коридор",
        ],
        conclusion="Проект обладает высокой перспективной транспортной доступностью.",
    ),
    Scene(
        name="education",
        title="ОБРАЗОВАНИЕ И ДОУ",
        bullets=[
            "Школы Ильинское-Усово",
            "Школы СберСити",
            "МГИМО",
            "Физтех-лицей",
            "Сколково (с 2012 г.)",
        ],
        conclusion="Локация формирует сильный образовательный кластер западного направления.",
    ),
    Scene(
        name="sport",
        title="СПОРТ И РЕКРЕАЦИЯ",
        bullets=[
            "Серебряный бор",
            "Мещерский парк",
            "Outdoor-инфраструктура",
            "Яхтинг и рекреация",
            "Лесные массивы",
        ],
        conclusion="Локация обладает выраженным природным и рекреационным потенциалом.",
    ),
    Scene(
        name="trade",
        title="ТОРГОВЛЯ",
        bullets=[
            "Vegas Crocus",
            "Твой Дом",
            "METRO Cash & Carry",
            "Глобус",
        ],
        conclusion="Крупный ритейл-кластер западного направления Москвы. Высокая транспортная доступность через Новорижское шоссе.",
    ),
    Scene(
        name="medical",
        title="МЕДИЦИНА",
        bullets=[
            "Клинический госпиталь Лапино",
            "Медицинский кластер Сколково",
            "Медицина СберСити",
        ],
        conclusion="Формирование медицинского и инновационного кластера. Доступ к высокотехнологичной медицине мирового уровня.",
        panel_blue=True,
    ),
]

# =============================================================================
# HELPERS
# =============================================================================

def clear_slide(slide) -> None:
    """Удаляет все shape'ы со слайда, не трогая сам slide part и отношения презентации."""
    sp_tree = slide.shapes._spTree
    for element in list(sp_tree):
        tag = element.tag
        # nvGrpSpPr / grpSpPr — служебные элементы группы shape tree, их трогать нельзя.
        if tag.endswith('}nvGrpSpPr') or tag.endswith('}grpSpPr'):
            continue
        sp_tree.remove(element)


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def move_last_slide_to_index(prs, target_index: int) -> None:
    """Перемещает только что добавленный последний слайд на target_index, не удаляя существующие слайды."""
    sld_id_lst = prs.slides._sldIdLst
    new_slide_id = sld_id_lst[-1]
    sld_id_lst.remove(new_slide_id)
    sld_id_lst.insert(target_index, new_slide_id)


def add_header(slide) -> None:
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.28))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.color.rgb = NAVY
    hdr.line.width = Pt(1)

    badge = slide.shapes.add_textbox(Inches(8.45), Inches(0.02), Inches(1.35), Inches(0.23))
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "bnMAP.pro"
    run.font.name = FONT_FAMILY
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_title(slide, title_text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.36), Inches(9.0), Inches(0.42))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT


def add_left_panel(slide, scene: Scene) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(1.02), Inches(3.45), Inches(3.62))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG_BLUE if scene.panel_blue else PANEL_BG
    panel.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
    panel.line.width = Pt(0.7)

    tf_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.18), Inches(3.05), Inches(3.25))
    tf = tf_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, bullet in enumerate(scene.bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = FONT_FAMILY
        run.font.size = Pt(11)
        run.font.color.rgb = DARK_TEXT


def add_map(slide, scene: Scene) -> None:
    map_path = os.path.join(MAPS_DIR, f"{scene.name}_map.png")
    if not os.path.exists(map_path):
        raise FileNotFoundError(f"Не найдена карта: {map_path}. Сначала запусти generate_maps_final.py")

    slide.shapes.add_picture(map_path, Inches(4.06), Inches(0.96), width=Inches(5.55), height=Inches(4.07))


def add_conclusion(slide, scene: Scene) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(4.78), Inches(9.18), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.62), Inches(4.86), Inches(8.78), Inches(0.39))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Основной вывод: " + scene.conclusion
    run.font.name = FONT_FAMILY
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE


def build_scene_slide(slide, scene: Scene) -> None:
    clear_slide(slide)
    add_header(slide)
    add_title(slide, scene.title)
    add_left_panel(slide, scene)
    add_map(slide, scene)
    add_conclusion(slide, scene)


def ensure_maps_exist() -> None:
    missing = []
    for scene in SCENES:
        p = os.path.join(MAPS_DIR, f"{scene.name}_map.png")
        if not os.path.exists(p):
            missing.append(p)
    if missing:
        joined = "\n".join(missing)
        raise FileNotFoundError("Не найдены карты:\n" + joined + "\nСначала запусти: python generate_maps_final.py")

# =============================================================================
# MAIN
# =============================================================================

def build() -> None:
    os.chdir(PROJECT_DIR)
    ensure_maps_exist()

    prs = Presentation(SOURCE_PPTX)

    if len(prs.slides) < 9:
        raise RuntimeError(f"В исходной презентации слишком мало слайдов: {len(prs.slides)}")

    # Слайды 4–8: location, transport, education, sport, trade.
    for slide_idx, scene in zip(range(3, 8), SCENES[:5]):
        build_scene_slide(prs.slides[slide_idx], scene)
        print(f"  ✓ слайд {slide_idx + 1}: {scene.title}")

    # Новый слайд 9: medicine.
    blank = get_blank_layout(prs)
    medical_slide = prs.slides.add_slide(blank)
    build_scene_slide(medical_slide, SCENES[5])
    move_last_slide_to_index(prs, 8)  # 0-based: девятая позиция
    print("  ✓ слайд 9: МЕДИЦИНА")

    # Слайд 15/конкуренты не трогаем: никаких map background overlay.
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ ГОТОВО: {OUTPUT_PPTX} ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()
