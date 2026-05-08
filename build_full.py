#!/usr/bin/env python3
"""
Сборка ТПУ Ильинская — v4.

Что делает:
- Перерисовывает слайды 4–8 (карты: местоположение, транспорт, образование, спорт, торговля).
- Добавляет слайд 9 — Медицина.
- Перестраивает слайд 15 — Конкурентное окружение (OSM-карта + левая панель).
- Исправляет слайды 11/12 (SWOT): убирает пустые bullet-строки (trailing \\n в <a:t>).
- Исправляет слайд 25: subtitle сдвигается вниз, не перекрывает разделительную полосу.
- Не трогает все остальные слайды.

Полный цикл генерации:
    python generate_maps_final_v2.py   # 7 карт → maps/
    python build_full.py               # сборка PPTX

Зависимости:
    pip install python-pptx pillow lxml
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# =============================================================================
# CONFIG
# =============================================================================

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
SOURCE_PPTX = os.path.join(PROJECT_DIR, "TPU_Ilinskaya_reworked_final.pptx")
OUTPUT_PPTX = os.path.join(PROJECT_DIR, "TPU_Ilinskaya_final_safe.pptx")
MAPS_DIR = os.path.join(PROJECT_DIR, "maps")

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.63)

# =============================================================================
# DESIGN SYSTEM
# =============================================================================

NAVY = RGBColor(0x0D, 0x21, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x19, 0x19, 0x19)
MUTED_TEXT = RGBColor(0x66, 0x77, 0x88)
PANEL_BG = RGBColor(0xF8, 0xF8, 0xF8)
PANEL_BLUE = RGBColor(0xED, 0xF4, 0xFA)
ACCENT = RGBColor(0x0B, 0x5D, 0x8C)
FONT = "Calibri"

# =============================================================================
# DATA — карточные слайды 4–9
# =============================================================================

@dataclass
class Scene:
    name: str
    title: str
    bullets: List[str]
    conclusion: str
    panel_blue: bool = False


SCENES: List[Scene] = [
    Scene(
        name="location",
        title="ОБ УЧАСТКЕ И МЕСТОПОЛОЖЕНИИ",
        bullets=[
            "Кадастровый номер: 50:11:0050603:423",
            "Площадь участка: ~20,6 га",
            "Форма участка: вытянутый линейный тип",
            "Север: Новорижское шоссе",
            "Юг: Москва-река",
            "Запад: Воронки / Архангельское",
            "Восток: Рублёво / МКАД направление",
        ],
        conclusion=(
            "Участок расположен в зоне активного развития западного направления Москвы."
        ),
    ),
    Scene(
        name="transport",
        title="ТРАНСПОРТНАЯ ДОСТУПНОСТЬ",
        bullets=[
            "Расстояние до МКАД: ~5 км",
            "Москва-Сити: ~40–55 мин",
            "м. Строгино: ~20–30 мин",
            "ст. Ильинская: шаговая доступность",
            "Новорижское шоссе — основной коридор",
        ],
        conclusion=(
            "Проект обладает высокой перспективной транспортной доступностью."
        ),
    ),
    Scene(
        name="education",
        title="ОБРАЗОВАНИЕ И ДОУ",
        bullets=[
            "Школы Ильинское-Усово",
            "Школы СберСити",
            "МГИМО",
            "Физтех-лицей",
            "Сколково (с 2012 г.)",
        ],
        conclusion=(
            "Локация формирует сильный образовательный кластер западного направления."
        ),
    ),
    Scene(
        name="sport",
        title="СПОРТ И РЕКРЕАЦИЯ",
        bullets=[
            "Серебряный бор",
            "Мещерский парк",
            "Outdoor-инфраструктура",
            "Яхтинг и рекреация",
            "Лесные массивы",
        ],
        conclusion=(
            "Локация обладает выраженным природным и рекреационным потенциалом."
        ),
    ),
    Scene(
        name="trade",
        title="ТОРГОВЛЯ",
        bullets=[
            "Vegas Crocus",
            "Твой Дом",
            "METRO Cash & Carry",
            "Глобус",
        ],
        conclusion=(
            "Крупный ритейл-кластер западного направления Москвы. "
            "Высокая транспортная доступность через Новорижское шоссе."
        ),
    ),
    Scene(
        name="medical",
        title="МЕДИЦИНА",
        bullets=[
            "Клинический госпиталь Лапино",
            "Медицинский кластер Сколково",
            "Медицина СберСити",
        ],
        conclusion=(
            "Формирование медицинского и инновационного кластера. "
            "Доступ к высокотехнологичной медицине мирового уровня."
        ),
        panel_blue=True,
    ),
]

# =============================================================================
# DATA — слайд 15: конкурентное окружение
# =============================================================================

COMPETITORS_BULLETS = [
    "Бизнес-класс (прямые конкуренты):",
    "  1 — ТПУ Ильинская (Целевой проект)",
    "  2 — СберСити (крупнейший, бренд СБ)",
    "  3 — Станиславский (прямой конкурент)",
    "  6 — Мыс (MR Group, Премиум/Бизнес)",
    "  7 — Резиденции Сколково",
    "Комфорт-класс (частичное пересечение ЦА):",
    "  4 — Строгино 360 (ПИК)",
    "  5 — Квартал Строгино (Самолёт)",
    "  8 — Сити Бэй (MR Group)",
    "  9 — Родина Переделкино",
]

COMPETITORS_CONCLUSION = (
    "Основную конкурентную угрозу в классе Бизнес представляет «СберСити» (масштаб, бренд). "
    "«Станиславский» — прямой конкурент по позиционированию. "
    "Проекты комфорт-класса частично пересекаются с 1 очередью Проекта."
)

# =============================================================================
# HELPERS — общий дизайн
# =============================================================================

def clear_slide(slide) -> None:
    """Удаляет все shape'ы со слайда, не трогая служебные элементы spTree."""
    sp_tree = slide.shapes._spTree
    for el in list(sp_tree):
        tag = el.tag
        if tag.endswith('}nvGrpSpPr') or tag.endswith('}grpSpPr'):
            continue
        sp_tree.remove(el)


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def move_last_slide_to_index(prs, target_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(target_index, new_id)


# =============================================================================
# HELPERS — элементы слайда
# =============================================================================

def add_header(slide) -> None:
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.28))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.color.rgb = NAVY
    hdr.line.width = Pt(1)

    badge = slide.shapes.add_textbox(
        Inches(8.45), Inches(0.02), Inches(1.35), Inches(0.23))
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "bnMAP.pro"
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_title(slide, title_text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.36), Inches(9.0), Inches(0.42))
    tf = box.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    run.font.name = FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT


def add_left_panel(slide, bullets: List[str], panel_blue: bool = False) -> None:
    """Левая панель: скруглённый прямоугольник + текст буллетов."""
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.42), Inches(1.02), Inches(3.45), Inches(3.62))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BLUE if panel_blue else PANEL_BG
    panel.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
    panel.line.width = Pt(0.7)

    tb = slide.shapes.add_textbox(
        Inches(0.62), Inches(1.12), Inches(3.05), Inches(3.35))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        # Категорийные заголовки (не отступают пробелами) — чуть крупнее и жирнее
        is_header = not bullet.startswith(" ")
        run.font.size = Pt(11) if is_header else Pt(10)
        run.font.bold = is_header
        run.font.color.rgb = DARK_TEXT


def add_map(slide, scene_name: str) -> None:
    map_path = os.path.join(MAPS_DIR, f"{scene_name}_map.png")
    if not os.path.exists(map_path):
        raise FileNotFoundError(
            f"Карта не найдена: {map_path}\n"
            "Сначала запусти: python generate_maps_final_v2.py"
        )
    slide.shapes.add_picture(
        map_path, Inches(4.06), Inches(0.96), width=Inches(5.55), height=Inches(4.07))


def add_competitors_map(slide) -> None:
    map_path = os.path.join(MAPS_DIR, "competitors_bg.png")
    if not os.path.exists(map_path):
        raise FileNotFoundError(
            f"Карта конкурентов не найдена: {map_path}\n"
            "Сначала запусти: python generate_maps_final_v2.py"
        )
    slide.shapes.add_picture(
        map_path, Inches(4.06), Inches(0.96), width=Inches(5.55), height=Inches(4.07))


def add_conclusion(slide, text: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.42), Inches(4.78), Inches(9.18), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(
        Inches(0.62), Inches(4.86), Inches(8.78), Inches(0.39))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = "Основной вывод: " + text
    run.font.name = FONT
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE


# =============================================================================
# ПОСТРОЕНИЕ СЛАЙДОВ
# =============================================================================

def build_scene_slide(slide, scene: Scene) -> None:
    """Перестраивает слайд по шаблону карты (слайды 4–9)."""
    clear_slide(slide)
    add_header(slide)
    add_title(slide, scene.title)
    add_left_panel(slide, [f"• {b}" for b in scene.bullets], scene.panel_blue)
    add_map(slide, scene.name)
    add_conclusion(slide, scene.conclusion)


def build_competitors_slide(slide) -> None:
    """Перестраивает слайд 15 — Конкурентное окружение в стиле карт-слайдов."""
    clear_slide(slide)
    add_header(slide)
    add_title(slide, "ПРОЕКТЫ КОНКУРЕНТНОГО ОКРУЖЕНИЯ В ЛОКАЦИИ")
    add_left_panel(slide, COMPETITORS_BULLETS, panel_blue=True)
    add_competitors_map(slide)
    add_conclusion(slide, COMPETITORS_CONCLUSION)


# =============================================================================
# ИСПРАВЛЕНИЯ СУЩЕСТВУЮЩИХ СЛАЙДОВ
# =============================================================================

def fix_swot_trailing_newlines(slide) -> int:
    """
    Убирает символ \\n в конце текста каждого <a:t> в слайде.
    Это устраняет лишние пустые bullet-строки в SWOT-слайдах.
    Возвращает количество исправленных элементов.
    """
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for t_el in shape.text_frame._txBody.iter(f'{{{NS_A}}}t'):
            if t_el.text and t_el.text.endswith('\n'):
                t_el.text = t_el.text.rstrip('\n')
                count += 1
    return count


def fix_slide25_subtitle(slide) -> bool:
    """
    Сдвигает subtitle-textbox на слайде 25 (Ценообразование на старте продаж)
    чтобы он не перекрывал разделительную синюю полосу.
    Возвращает True если элемент найден и исправлен.
    """
    keywords = ('методика', 'предпосылки', 'редика', 'редпосылки')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip().lower()
        if any(kw in txt for kw in keywords):
            shape.top = Inches(0.62)
            return True
    return False


# =============================================================================
# ПРОВЕРКА НАЛИЧИЯ КАРТ
# =============================================================================

def ensure_maps_exist() -> None:
    missing = []
    for scene in SCENES:
        p = os.path.join(MAPS_DIR, f"{scene.name}_map.png")
        if not os.path.exists(p):
            missing.append(p)
    comp = os.path.join(MAPS_DIR, "competitors_bg.png")
    if not os.path.exists(comp):
        missing.append(comp)
    if missing:
        raise FileNotFoundError(
            "Не найдены карты:\n" + "\n".join(missing) +
            "\n\nСначала запусти: python generate_maps_final_v2.py"
        )


# =============================================================================
# MAIN
# =============================================================================

def build() -> None:
    os.chdir(PROJECT_DIR)
    ensure_maps_exist()

    prs = Presentation(SOURCE_PPTX)

    if len(prs.slides) < 15:
        raise RuntimeError(
            f"В исходной презентации слишком мало слайдов: {len(prs.slides)}"
        )

    # ── Слайды 4–8: карты местоположение–торговля ──────────────────────────
    for slide_idx, scene in zip(range(3, 8), SCENES[:5]):
        build_scene_slide(prs.slides[slide_idx], scene)
        print(f"  ✓ слайд {slide_idx + 1}: {scene.title}")

    # ── Новый слайд 9: медицина ─────────────────────────────────────────────
    blank = get_blank_layout(prs)
    medical_slide = prs.slides.add_slide(blank)
    build_scene_slide(medical_slide, SCENES[5])
    move_last_slide_to_index(prs, 8)  # 0-based → позиция 9
    print("  ✓ слайд 9: МЕДИЦИНА")

    # ── Слайд 15: конкурентное окружение ───────────────────────────────────
    # После добавления слайда 9 все следующие сдвинулись на +1,
    # поэтому слайд 15 теперь находится на индексе 15 (было 14).
    competitors_idx = 15
    build_competitors_slide(prs.slides[competitors_idx])
    print("  ✓ слайд 16 (бывший 15): КОНКУРЕНТНОЕ ОКРУЖЕНИЕ — OSM-карта")

    # ── Исправление SWOT слайдов 11/12 → после вставки слайда 9: 11/12 ─────
    for idx, label in [(11, "11"), (12, "12")]:
        n = fix_swot_trailing_newlines(prs.slides[idx])
        print(f"  ✓ слайд {idx + 1} (SWOT): исправлено {n} trailing \\n")

    # ── Исправление слайда 25 (после вставки = индекс 25) ──────────────────
    fixed = fix_slide25_subtitle(prs.slides[25])
    print(f"  ✓ слайд 26: subtitle {'сдвинут' if fixed else 'не найден — проверь вручную'}")

    prs.save(OUTPUT_PPTX)
    print(f"\n✅ ГОТОВО: {OUTPUT_PPTX} ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()
